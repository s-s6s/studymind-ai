import fitz  # PyMuPDF
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import io


class FileProcessor:
    async def extract_text(self, file_bytes: bytes, file_type: str) -> str:
        extractors = {
            "pdf": self._extract_pdf,
            "docx": self._extract_docx,
            "pptx": self._extract_pptx,
            "image": self._extract_image,
            "txt": self._extract_txt,
        }
        extractor = extractors.get(file_type)
        if not extractor:
            raise ValueError(f"Unsupported file type: {file_type}")
        return await extractor(file_bytes)

    async def _extract_pdf(self, file_bytes: bytes) -> str:
        text = ""
        # Try PyMuPDF first
        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
        except Exception:
            pass

        # Fallback to pdfplumber
        if len(text.strip()) < 100:
            try:
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except Exception:
                pass

        # Fallback to OCR for scanned PDFs
        if len(text.strip()) < 100:
            try:
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                for page in doc:
                    pix = page.get_pixmap(dpi=300)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text += pytesseract.image_to_string(img, lang='ara+eng') + "\n"
                doc.close()
            except Exception:
                pass

        return text.strip()

    async def _extract_docx(self, file_bytes: bytes) -> str:
        doc = Document(io.BytesIO(file_bytes))
        text = ""
        for p in doc.paragraphs:
            text += p.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        return text.strip()

    async def _extract_pptx(self, file_bytes: bytes) -> str:
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {i} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text.strip()

    async def _extract_image(self, file_bytes: bytes) -> str:
        img = Image.open(io.BytesIO(file_bytes))
        return pytesseract.image_to_string(img, lang='ara+eng').strip()

    async def _extract_txt(self, file_bytes: bytes) -> str:
        for enc in ['utf-8', 'utf-8-sig', 'windows-1256', 'iso-8859-6']:
            try:
                return file_bytes.decode(enc).strip()
            except Exception:
                continue
        return file_bytes.decode('utf-8', errors='ignore').strip()

    def detect_file_type(self, filename: str, content_type: str) -> str:
        ext = filename.lower().rsplit('.', 1)[-1] if '.' in filename else ''
        type_map = {
            'pdf': 'pdf', 'docx': 'docx', 'doc': 'docx',
            'pptx': 'pptx', 'ppt': 'pptx',
            'png': 'image', 'jpg': 'image', 'jpeg': 'image',
            'txt': 'txt',
        }
        return type_map.get(ext, 'txt')
